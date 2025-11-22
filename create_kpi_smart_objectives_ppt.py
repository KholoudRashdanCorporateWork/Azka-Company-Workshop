#!/usr/bin/env python3
"""
Create a comprehensive PowerPoint presentation on KPIs and SMART Objectives
for a 2-day workshop targeting middle managers in the technology sector.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION


def set_text_format(text_frame, font_name="Calibri", font_size=18, bold=False, color=RGBColor(0, 0, 0)):
    """Helper function to format text"""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 114, 196)

    return slide


def add_content_slide(prs, title, content_list):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for item in content_list:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.name = "Calibri"
        p.space_after = Pt(10)

    return slide


def add_two_column_slide(prs, title, left_content, right_content):
    """Add a slide with two columns"""
    slide_layout = prs.slide_layouts[3]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Left column
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(5)
    left_box = slide.shapes.add_textbox(left, top, width, height)
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    # Right column
    left = Inches(5.2)
    right_box = slide.shapes.add_textbox(left, top, width, height)
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    return slide


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Add table
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    num_rows = len(rows) + 1
    num_cols = len(headers)

    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Set headers
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True
                run.font.size = Pt(16)

    # Set rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(220, 230, 241)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)

    return slide


def add_chart_slide(prs, title, chart_data_dict, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED):
    """Add a slide with a chart"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Add chart
    chart_data = CategoryChartData()
    chart_data.categories = chart_data_dict['categories']

    for series_name, values in chart_data_dict['series'].items():
        chart_data.add_series(series_name, values)

    x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(4.5)
    chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    return slide


def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== TITLE SLIDE ==========
    add_title_slide(
        prs,
        "How to Write Effective KPIs and SMART Objectives",
        "A 2-Day Workshop for Technology Middle Managers"
    )

    # ========== WORKSHOP OVERVIEW ==========
    add_content_slide(prs, "Workshop Overview", [
        "Day 1: Understanding SMART Objectives & Goal Setting",
        "• What are objectives and why they matter",
        "• Introduction to the SMART framework",
        "• Writing effective function/team objectives",
        "• Practical exercises and examples",
        "",
        "Day 2: KPIs & Cascading to Your Team",
        "• Understanding KPIs and their importance",
        "• Creating meaningful KPIs from objectives",
        "• Cascading objectives and KPIs to team members",
        "• Real-world technology sector examples"
    ])

    # ========== LEARNING OUTCOMES ==========
    add_content_slide(prs, "Learning Outcomes", [
        "By the end of this workshop, you will be able to:",
        "",
        "✓ Understand the difference between objectives and KPIs",
        "✓ Write SMART objectives for your function/team",
        "✓ Create measurable KPIs that drive performance",
        "✓ Cascade organizational objectives to your team",
        "✓ Align individual KPIs with team and company goals",
        "✓ Monitor and track performance effectively"
    ])

    # ========== DAY 1: SECTION DIVIDER ==========
    add_title_slide(prs, "DAY 1", "Understanding SMART Objectives & Goal Setting")

    # ========== WHAT ARE OBJECTIVES ==========
    add_content_slide(prs, "What Are Objectives?", [
        "Definition:",
        "Objectives are specific, measurable goals that define what you want to achieve",
        "",
        "Why Objectives Matter:",
        "• Provide clear direction and focus",
        "• Align team efforts with organizational strategy",
        "• Enable performance measurement",
        "• Motivate and engage team members",
        "• Facilitate resource allocation and prioritization"
    ])

    # ========== OBJECTIVES VS KPIs ==========
    add_table_slide(prs, "Objectives vs KPIs: Understanding the Difference",
        ["Aspect", "Objectives", "KPIs"],
        [
            ["Definition", "What you want to achieve", "How you measure achievement"],
            ["Nature", "Qualitative or Quantitative", "Always Quantitative"],
            ["Purpose", "Set direction and goals", "Track progress and performance"],
            ["Example", "Improve customer satisfaction", "NPS score of 8.5 or higher"],
            ["Time Frame", "Medium to long-term", "Measured regularly (daily/weekly/monthly)"],
            ["Focus", "Outcome-oriented", "Metric-oriented"]
        ]
    )

    # ========== THE SMART FRAMEWORK ==========
    add_title_slide(prs, "The SMART Framework", "A Proven Method for Effective Objectives")

    # ========== SMART - SPECIFIC ==========
    add_two_column_slide(prs, "S - Specific",
        [
            "What does SPECIFIC mean?",
            "• Clearly defined and unambiguous",
            "• Answers: Who, What, Where, When, Why",
            "• Leaves no room for misinterpretation",
            "• Focuses on a single objective",
            "",
            "Why it matters:",
            "• Provides clarity and direction",
            "• Reduces confusion",
            "• Easier to communicate to team"
        ],
        [
            "❌ Poor Example:",
            "'Improve our software quality'",
            "",
            "✅ SMART Example:",
            "'Reduce critical bugs in production by implementing automated testing for all new features in our mobile application'",
            "",
            "Notice the difference:",
            "• What: Reduce critical bugs",
            "• How: Automated testing",
            "• Where: Mobile application"
        ]
    )

    # ========== SMART - MEASURABLE ==========
    add_two_column_slide(prs, "M - Measurable",
        [
            "What does MEASURABLE mean?",
            "• Can be quantified or qualified",
            "• Has clear criteria for success",
            "• Progress can be tracked",
            "• Includes specific numbers/percentages",
            "",
            "Why it matters:",
            "• Enables progress tracking",
            "• Provides objective evaluation",
            "• Motivates through visible progress"
        ],
        [
            "❌ Poor Example:",
            "'Make our development process faster'",
            "",
            "✅ SMART Example:",
            "'Reduce average sprint cycle time from 3 weeks to 2 weeks by Q2 2024'",
            "",
            "Measurement criteria:",
            "• Baseline: 3 weeks",
            "• Target: 2 weeks",
            "• Metric: Sprint cycle time",
            "• Improvement: 33% reduction"
        ]
    )

    # ========== SMART - ACHIEVABLE ==========
    add_two_column_slide(prs, "A - Achievable",
        [
            "What does ACHIEVABLE mean?",
            "• Realistic given resources",
            "• Challenging but attainable",
            "• Within your control/influence",
            "• Considers constraints and risks",
            "",
            "Why it matters:",
            "• Maintains team motivation",
            "• Prevents burnout",
            "• Builds credibility",
            "• Ensures sustainable progress"
        ],
        [
            "❌ Poor Example:",
            "'Achieve 100% code coverage across all 50 legacy applications in 1 month'",
            "",
            "✅ SMART Example:",
            "'Achieve 80% code coverage on 5 priority applications within 6 months by allocating 20% of sprint capacity'",
            "",
            "Why it's achievable:",
            "• Realistic scope (5 apps, not 50)",
            "• Reasonable timeline (6 months)",
            "• Resource allocation defined (20%)",
            "• Practical target (80% vs 100%)"
        ]
    )

    # ========== SMART - RELEVANT ==========
    add_two_column_slide(prs, "R - Relevant",
        [
            "What does RELEVANT mean?",
            "• Aligns with broader goals",
            "• Supports organizational strategy",
            "• Matters to stakeholders",
            "• Worth the time and effort",
            "",
            "Why it matters:",
            "• Ensures strategic alignment",
            "• Maximizes impact",
            "• Justifies resource investment",
            "• Maintains focus on priorities"
        ],
        [
            "Context: Company goal is to improve customer retention by 25%",
            "",
            "❌ Poor Example (Not Relevant):",
            "'Migrate all internal tools to newest framework version'",
            "",
            "✅ SMART Example (Relevant):",
            "'Reduce customer-reported app crashes by 60% to improve retention and user satisfaction'",
            "",
            "Why it's relevant:",
            "• Directly impacts customer experience",
            "• Supports retention goal",
            "• Addresses customer pain point"
        ]
    )

    # ========== SMART - TIME-BOUND ==========
    add_two_column_slide(prs, "T - Time-Bound",
        [
            "What does TIME-BOUND mean?",
            "• Has a clear deadline",
            "• Includes milestones",
            "• Creates urgency",
            "• Enables progress tracking",
            "",
            "Why it matters:",
            "• Prevents procrastination",
            "• Enables planning and scheduling",
            "• Creates accountability",
            "• Allows for course correction"
        ],
        [
            "❌ Poor Example:",
            "'Implement new CI/CD pipeline eventually'",
            "",
            "✅ SMART Example:",
            "'Implement CI/CD pipeline for 3 core services by end of Q3 2024, with Phase 1 (design) by July 15 and Phase 2 (implementation) by Sept 15'",
            "",
            "Timeline breakdown:",
            "• Phase 1 deadline: July 15, 2024",
            "• Phase 2 deadline: Sept 15, 2024",
            "• Final completion: Q3 2024",
            "• Scope: 3 core services"
        ]
    )

    # ========== SMART EXAMPLES COMPARISON ==========
    add_table_slide(prs, "Before & After: SMART Transformation",
        ["Before (Weak)", "After (SMART)", "What Changed"],
        [
            [
                "Improve team skills",
                "Complete cloud certification for 8 team members by Q4 2024 with 90% pass rate",
                "Added specific number, metric, and deadline"
            ],
            [
                "Better customer support",
                "Reduce average ticket resolution time from 48h to 24h by implementing new ticketing system by July 2024",
                "Quantified improvement, added method and timeline"
            ],
            [
                "Modernize infrastructure",
                "Migrate 75% of on-premise workloads to cloud by Dec 2024, starting with 3 pilot applications in Q2",
                "Specified percentage, deadline, and phased approach"
            ],
            [
                "Increase productivity",
                "Reduce manual deployment time by 70% (from 4h to 1.2h) by automating release process by Q3 2024",
                "Quantified baseline, target, and method"
            ]
        ]
    )

    # ========== WRITING FUNCTION OBJECTIVES ==========
    add_content_slide(prs, "Writing Your Function/Team Objectives", [
        "Step 1: Understand organizational strategy and goals",
        "• Review company OKRs, strategic plans, and priorities",
        "• Identify how your function contributes",
        "",
        "Step 2: Identify your function's key responsibilities",
        "• What are you accountable for?",
        "• What value does your team provide?",
        "",
        "Step 3: Define 3-5 key objectives (don't overcommit!)",
        "• Focus on high-impact areas",
        "• Apply SMART criteria to each",
        "",
        "Step 4: Validate with stakeholders",
        "• Get buy-in from leadership and peers"
    ])

    # ========== TECHNOLOGY FUNCTION EXAMPLES ==========
    add_table_slide(prs, "Sample Objectives by Technology Function",
        ["Function", "Sample SMART Objective"],
        [
            [
                "Software Development",
                "Increase deployment frequency from monthly to weekly by implementing automated CI/CD pipeline by Q3 2024"
            ],
            [
                "Infrastructure & DevOps",
                "Achieve 99.9% system uptime by implementing redundancy and automated failover for all critical services by Q4 2024"
            ],
            [
                "Quality Assurance",
                "Reduce production defects by 50% (from 20 to 10 per release) by implementing shift-left testing by Q2 2024"
            ],
            [
                "IT Security",
                "Complete security compliance audit with zero critical findings by implementing automated vulnerability scanning by Dec 2024"
            ],
            [
                "Product Management",
                "Increase feature adoption rate from 35% to 60% by implementing user feedback loop and A/B testing by Q3 2024"
            ],
            [
                "Data Engineering",
                "Reduce data processing latency from 24h to 2h by optimizing ETL pipelines for 5 core data sources by Q4 2024"
            ]
        ]
    )

    # ========== COMMON PITFALLS ==========
    add_table_slide(prs, "Common Pitfalls to Avoid",
        ["Pitfall", "Why It's a Problem", "How to Fix It"],
        [
            [
                "Too many objectives",
                "Dilutes focus and resources",
                "Limit to 3-5 key objectives"
            ],
            [
                "Vague language",
                "Open to interpretation",
                "Use specific, measurable terms"
            ],
            [
                "No baseline data",
                "Can't measure improvement",
                "Establish current state first"
            ],
            [
                "Unrealistic targets",
                "Demotivates team",
                "Set challenging but achievable goals"
            ],
            [
                "Activity-focused",
                "Measures effort, not outcome",
                "Focus on results and impact"
            ],
            [
                "Missing stakeholder buy-in",
                "Misalignment and conflicts",
                "Validate with leadership early"
            ]
        ]
    )

    # ========== DAY 1 EXERCISE ==========
    add_content_slide(prs, "🎯 Day 1 Exercise: Write Your Own Objectives", [
        "Activity (45 minutes):",
        "",
        "1. Review your function's current responsibilities (10 min)",
        "   • List your team's key accountabilities",
        "",
        "2. Draft 3-5 SMART objectives for your function (25 min)",
        "   • Use the SMART framework",
        "   • Consider organizational alignment",
        "",
        "3. Peer review in pairs (10 min)",
        "   • Exchange with a colleague",
        "   • Check each criterion: S-M-A-R-T",
        "   • Provide constructive feedback",
        "",
        "We'll share examples after the break!"
    ])

    # ========== DAY 1 SUMMARY ==========
    add_content_slide(prs, "Day 1 Summary", [
        "Key Takeaways:",
        "",
        "✓ Objectives define WHAT you want to achieve",
        "✓ SMART framework ensures objectives are effective:",
        "  • Specific, Measurable, Achievable, Relevant, Time-bound",
        "✓ Good objectives provide clarity and direction",
        "✓ Function objectives must align with organizational strategy",
        "✓ Quality over quantity - focus on 3-5 key objectives",
        "",
        "Tomorrow: We'll learn how to create KPIs and cascade them to your team members!"
    ])

    # ========== DAY 2: SECTION DIVIDER ==========
    add_title_slide(prs, "DAY 2", "KPIs & Cascading to Your Team")

    # ========== WHAT ARE KPIs ==========
    add_content_slide(prs, "What Are Key Performance Indicators (KPIs)?", [
        "Definition:",
        "KPIs are quantifiable metrics that measure progress toward objectives",
        "",
        "Characteristics of Good KPIs:",
        "• Quantifiable - expressed in numbers or percentages",
        "• Actionable - can influence through your actions",
        "• Relevant - directly related to objectives",
        "• Timely - measured at appropriate intervals",
        "• Accurate - based on reliable data",
        "",
        "Remember: KPIs measure HOW you track achievement of objectives"
    ])

    # ========== TYPES OF KPIs ==========
    add_table_slide(prs, "Types of KPIs in Technology",
        ["KPI Type", "What It Measures", "Example"],
        [
            [
                "Input KPIs",
                "Resources consumed",
                "Hours spent on development, Budget allocated"
            ],
            [
                "Process KPIs",
                "Efficiency of activities",
                "Code review turnaround time, Sprint velocity"
            ],
            [
                "Output KPIs",
                "Deliverables produced",
                "Features deployed, Stories completed"
            ],
            [
                "Outcome KPIs",
                "Results achieved",
                "System uptime %, Customer satisfaction score"
            ],
            [
                "Leading KPIs",
                "Predict future performance",
                "Test coverage %, Number of code commits"
            ],
            [
                "Lagging KPIs",
                "Past performance",
                "Production incidents, Customer churn rate"
            ]
        ]
    )

    # ========== FROM OBJECTIVES TO KPIs ==========
    add_content_slide(prs, "How to Create KPIs from Objectives", [
        "Step 1: Start with your SMART objective",
        "• The objective already contains measurable elements",
        "",
        "Step 2: Identify what needs to be measured",
        "• What indicates progress?",
        "• What shows success?",
        "",
        "Step 3: Define the metric formula",
        "• How will you calculate it?",
        "• What data sources will you use?",
        "",
        "Step 4: Set target values and thresholds",
        "• What's the goal? (Target)",
        "• What's acceptable? (Threshold)",
        "• What's excellent? (Stretch goal)"
    ])

    # ========== KPI EXAMPLES FROM OBJECTIVES ==========
    add_table_slide(prs, "From Objectives to KPIs: Examples",
        ["SMART Objective", "Related KPIs", "Target"],
        [
            [
                "Reduce production bugs by 50% by Q4 2024",
                "• # of critical bugs per release\n• Bug resolution time (hours)\n• Bug escape rate (%)",
                "• ≤ 5 critical bugs\n• < 24 hours\n• < 5%"
            ],
            [
                "Increase deployment frequency from monthly to weekly by Q3",
                "• Deployments per week\n• Deployment success rate (%)\n• Mean time to deploy (min)",
                "• 4 per week\n• ≥ 95%\n• < 30 min"
            ],
            [
                "Improve system uptime to 99.9% by implementing redundancy by Q4",
                "• System uptime (%)\n• Mean time between failures (hours)\n• Mean time to recovery (min)",
                "• ≥ 99.9%\n• > 720 hours\n• < 15 min"
            ],
            [
                "Achieve 80% test automation coverage for critical features by Q2",
                "• Test automation coverage (%)\n• Test execution time (min)\n• Tests passing rate (%)",
                "• ≥ 80%\n• < 20 min\n• ≥ 98%"
            ]
        ]
    )

    # ========== KPI CHARACTERISTICS ==========
    add_two_column_slide(prs, "Characteristics of Effective KPIs",
        [
            "Good KPIs Are:",
            "",
            "✓ Clearly Defined",
            "• No ambiguity in calculation",
            "• Everyone understands the metric",
            "",
            "✓ Easy to Measure",
            "• Data is readily available",
            "• Can be tracked regularly",
            "",
            "✓ Actionable",
            "• Team can influence the outcome",
            "• Clear actions drive improvement"
        ],
        [
            "✓ Relevant",
            "• Directly tied to objectives",
            "• Matters to stakeholders",
            "",
            "✓ Balanced",
            "• Mix of leading & lagging indicators",
            "• Covers quality, speed, and efficiency",
            "",
            "✓ Limited in Number",
            "• 3-5 KPIs per objective",
            "• Focus on what matters most",
            "",
            "✓ Regularly Reviewed",
            "• Tracked at appropriate frequency",
            "• Used for decision-making"
        ]
    )

    # ========== KPI DASHBOARD EXAMPLE ==========
    chart_data = {
        'categories': ['Q1', 'Q2', 'Q3', 'Q4'],
        'series': {
            'Deployment Frequency': [4, 8, 12, 16],
            'Target': [16, 16, 16, 16]
        }
    }
    add_chart_slide(prs, "KPI Tracking Example: Deployment Frequency", chart_data, XL_CHART_TYPE.LINE)

    # ========== CASCADING OBJECTIVES ==========
    add_title_slide(prs, "Cascading Objectives & KPIs", "From Organization to Individual")

    # ========== WHY CASCADE ==========
    add_content_slide(prs, "Why Cascade Objectives and KPIs?", [
        "Benefits of Cascading:",
        "",
        "✓ Alignment - Everyone works toward common goals",
        "✓ Clarity - Each person understands their contribution",
        "✓ Accountability - Clear ownership at every level",
        "✓ Motivation - People see how their work matters",
        "✓ Transparency - Visible connection from top to bottom",
        "✓ Efficiency - Eliminates conflicting priorities",
        "",
        "The Goal: Create a 'line of sight' from company strategy to daily tasks"
    ])

    # ========== CASCADING FRAMEWORK ==========
    add_table_slide(prs, "The Cascading Framework: Levels of Objectives",
        ["Level", "Scope", "Example", "Owner"],
        [
            [
                "Organizational",
                "Company-wide strategic goals",
                "Become market leader in cloud solutions",
                "CEO / Executive Team"
            ],
            [
                "Departmental",
                "Division or department goals",
                "Achieve 99.9% platform reliability",
                "CTO / VP Engineering"
            ],
            [
                "Function/Team",
                "Team or function-specific goals",
                "Reduce infrastructure incidents by 60%",
                "Engineering Manager"
            ],
            [
                "Individual",
                "Personal performance goals",
                "Implement monitoring for 10 critical services",
                "DevOps Engineer"
            ]
        ]
    )

    # ========== CASCADE PROCESS ==========
    add_content_slide(prs, "The 5-Step Cascading Process", [
        "Step 1: Understand Upper-Level Objectives",
        "• Review organizational and departmental objectives",
        "• Identify how your function contributes",
        "",
        "Step 2: Define Your Function/Team Objectives",
        "• What must your team achieve to support upper objectives?",
        "• Apply SMART framework",
        "",
        "Step 3: Break Down into Individual Objectives",
        "• What should each team member accomplish?",
        "• Ensure fair distribution and capabilities match",
        "",
        "Step 4: Create Corresponding KPIs at Each Level",
        "• Team KPIs roll up to function objectives",
        "• Individual KPIs roll up to team objectives",
        "",
        "Step 5: Review and Align",
        "• Validate with team members",
        "• Ensure understanding and buy-in"
    ])

    # ========== CASCADE EXAMPLE ==========
    add_table_slide(prs, "Cascading Example: From Company to Individual",
        ["Level", "Objective", "KPI"],
        [
            [
                "Company",
                "Increase customer retention by 25% by end of year",
                "Customer retention rate ≥ 85%"
            ],
            [
                "Department (Engineering)",
                "Improve platform stability and performance by Q4",
                "System uptime ≥ 99.9%\nP95 response time < 200ms"
            ],
            [
                "Team (DevOps)",
                "Reduce production incidents by 60% by Q3 2024",
                "Incidents per month ≤ 4\nMTTR < 30 minutes"
            ],
            [
                "Individual (DevOps Engineer)",
                "Implement monitoring and alerting for 10 critical services by July 2024",
                "Services monitored: 10/10\nAlert response time < 5 min\nFalse positive rate < 5%"
            ]
        ]
    )

    # ========== INDIVIDUAL KPI DEVELOPMENT ==========
    add_content_slide(prs, "Creating Individual Team Member KPIs", [
        "Principles for Individual KPIs:",
        "",
        "• Based on role and responsibilities",
        "  → Software Engineer: Code quality, delivery, collaboration",
        "  → QA Engineer: Test coverage, defect detection, automation",
        "",
        "• Aligned with team objectives",
        "  → Individual KPIs contribute to team KPIs",
        "",
        "• Within individual's control",
        "  → Can be influenced by their actions",
        "",
        "• Balanced scorecard approach",
        "  → Mix of quantitative and qualitative metrics",
        "  → Include both results and behaviors",
        "",
        "• Developed collaboratively",
        "  → Discuss with team member, don't impose"
    ])

    # ========== INDIVIDUAL KPI EXAMPLES ==========
    add_table_slide(prs, "Individual KPI Examples by Role",
        ["Role", "Sample KPIs", "Target"],
        [
            [
                "Software Engineer",
                "• Story points delivered per sprint\n• Code review turnaround time\n• Code quality score (SonarQube)",
                "• 20-25 points\n• < 24 hours\n• A rating"
            ],
            [
                "QA Engineer",
                "• Test cases automated per sprint\n• Defect detection rate\n• Test execution time reduction",
                "• 15-20 cases\n• ≥ 85%\n• 30% reduction"
            ],
            [
                "DevOps Engineer",
                "• Pipeline reliability (%)\n• Infrastructure as Code coverage\n• Deployment automation %",
                "• ≥ 98%\n• ≥ 80%\n• ≥ 90%"
            ],
            [
                "Tech Lead",
                "• Team velocity improvement\n• Technical debt reduction\n• Knowledge sharing sessions led",
                "• 15% increase\n• 20% reduction\n• 2 per month"
            ],
            [
                "Product Owner",
                "• Feature adoption rate\n• Sprint goal achievement\n• Stakeholder satisfaction score",
                "• ≥ 60%\n• ≥ 90%\n• ≥ 8/10"
            ]
        ]
    )

    # ========== ALIGNMENT CHECK ==========
    add_content_slide(prs, "Ensuring Alignment: The Vertical Check", [
        "Ask These Questions:",
        "",
        "Bottom-Up (Individual → Company):",
        "• 'If I achieve my KPIs, does it help my team achieve theirs?'",
        "• 'If our team achieves our KPIs, does it help the department?'",
        "• 'Does the department's success support company objectives?'",
        "",
        "Top-Down (Company → Individual):",
        "• 'How does company strategy translate to department goals?'",
        "• 'What must my team do to support department objectives?'",
        "• 'What should each individual contribute?'",
        "",
        "If the answer to any question is unclear, revisit the cascade!"
    ])

    # ========== COMMON CASCADING MISTAKES ==========
    add_table_slide(prs, "Common Cascading Mistakes to Avoid",
        ["Mistake", "Impact", "Solution"],
        [
            [
                "Cascading too many objectives",
                "Team overwhelmed, diluted focus",
                "Limit to 3-5 objectives per level"
            ],
            [
                "Lost in translation",
                "Individual KPIs don't support team goals",
                "Verify alignment at each level"
            ],
            [
                "One-way communication",
                "Lack of buy-in, unrealistic targets",
                "Collaborate with team on objectives"
            ],
            [
                "Conflicting KPIs",
                "Competing priorities, confusion",
                "Review for contradictions before finalizing"
            ],
            [
                "No visibility",
                "Team doesn't see the big picture",
                "Share company and department objectives"
            ],
            [
                "Set and forget",
                "Objectives become irrelevant",
                "Review and adjust quarterly"
            ]
        ]
    )

    # ========== TRACKING AND MONITORING ==========
    add_content_slide(prs, "Tracking and Monitoring KPIs", [
        "Best Practices:",
        "",
        "• Establish measurement cadence",
        "  → Daily: Operational metrics (uptime, incidents)",
        "  → Weekly: Team performance (velocity, completed stories)",
        "  → Monthly: Strategic metrics (customer satisfaction, retention)",
        "",
        "• Use dashboards and visualization",
        "  → Make KPIs visible to the team",
        "  → Use tools like Jira, Grafana, Tableau",
        "",
        "• Regular review meetings",
        "  → Weekly team syncs, Monthly business reviews",
        "",
        "• Take action on insights",
        "  → KPIs should drive decisions, not just report status",
        "  → If trending wrong direction, investigate and adjust"
    ])

    # ========== FEEDBACK AND ADJUSTMENT ==========
    add_content_slide(prs, "Regular Reviews and Adjustments", [
        "Quarterly Review Process:",
        "",
        "1. Review actual performance vs. targets",
        "   • What went well? What didn't?",
        "",
        "2. Identify obstacles and enablers",
        "   • What helped or hindered progress?",
        "",
        "3. Celebrate achievements",
        "   • Recognize team and individual wins",
        "",
        "4. Adjust objectives and KPIs if needed",
        "   • Business context changes - objectives should too",
        "   • But don't change too frequently (stability matters)",
        "",
        "5. Plan for next quarter",
        "   • Set new targets, identify needed support"
    ])

    # ========== COMMUNICATION TIPS ==========
    add_content_slide(prs, "Communicating Objectives and KPIs to Your Team", [
        "Effective Communication Strategies:",
        "",
        "✓ Explain the 'Why'",
        "  • Connect objectives to company strategy and purpose",
        "",
        "✓ Be Transparent",
        "  • Share both successes and challenges",
        "",
        "✓ Encourage Questions",
        "  • Create safe space for clarification and feedback",
        "",
        "✓ Make it Visual",
        "  • Use dashboards, charts, and progress indicators",
        "",
        "✓ Regular Updates",
        "  • Don't wait for formal reviews - communicate continuously",
        "",
        "✓ Two-Way Dialogue",
        "  • Listen to team input on feasibility and approach"
    ])

    # ========== DAY 2 EXERCISE ==========
    add_content_slide(prs, "🎯 Day 2 Exercise: Cascade Your Objectives", [
        "Activity (60 minutes):",
        "",
        "Part 1: Create Team KPIs (20 min)",
        "• Take one of your function objectives from Day 1",
        "• Define 3-5 KPIs that measure progress",
        "• Specify targets and measurement frequency",
        "",
        "Part 2: Cascade to Individual (25 min)",
        "• Select 2 team members (or use fictional examples)",
        "• Define individual objectives that support team objective",
        "• Create 3-4 individual KPIs for each",
        "",
        "Part 3: Alignment Check (15 min)",
        "• Verify individual KPIs → Team KPIs → Function Objective",
        "• Share with a peer for feedback"
    ])

    # ========== REAL WORLD EXAMPLE ==========
    add_table_slide(prs, "Complete Example: E-Commerce Platform Team",
        ["Level", "Objective/KPI", "Owner"],
        [
            [
                "Company Objective",
                "Increase annual revenue by 30% through improved digital experience",
                "CEO"
            ],
            [
                "Engineering Dept KPI",
                "Improve platform performance: Page load time < 2 sec, 99.95% uptime",
                "VP Engineering"
            ],
            [
                "Platform Team Objective",
                "Optimize checkout process to reduce cart abandonment by 25% by Q3",
                "Engineering Manager"
            ],
            [
                "Team KPIs",
                "• Cart abandonment rate < 15%\n• Checkout completion time < 45 sec\n• Payment success rate > 99%",
                "Platform Team"
            ],
            [
                "Individual (Frontend Dev)",
                "Implement one-click checkout for returning users by July 2024",
                "Sarah Chen"
            ],
            [
                "Individual KPIs",
                "• Feature completion: 100% by July 31\n• Page load time: < 1.5 sec\n• Zero critical bugs in production",
                "Sarah Chen"
            ]
        ]
    )

    # ========== TOOLS AND TEMPLATES ==========
    add_content_slide(prs, "Tools and Templates", [
        "Recommended Tools for KPI Tracking:",
        "",
        "• Project Management: Jira, Azure DevOps, Monday.com",
        "• Dashboards: Grafana, Tableau, Power BI, Klipfolio",
        "• OKR Software: Weekdone, Perdoo, Gtmhub, WorkBoard",
        "• Custom: Google Sheets, Excel with regular updates",
        "",
        "Templates to Use:",
        "• Objective Setting Template (SMART format)",
        "• KPI Definition Sheet (metric, formula, target, frequency)",
        "• Cascading Matrix (showing alignment across levels)",
        "• Progress Tracker (weekly/monthly updates)",
        "• Review Meeting Agenda"
    ])

    # ========== OBJECTIVE TEMPLATE ==========
    add_table_slide(prs, "Template: SMART Objective Definition",
        ["Component", "Your Response"],
        [
            ["Specific: What exactly will you achieve?", ""],
            ["Measurable: How will you measure success?", ""],
            ["Achievable: Why is this realistic?", ""],
            ["Relevant: How does this align with broader goals?", ""],
            ["Time-bound: What is the deadline/timeline?", ""],
            ["Complete SMART Objective:", ""],
        ]
    )

    # ========== KPI TEMPLATE ==========
    add_table_slide(prs, "Template: KPI Definition Sheet",
        ["Element", "Description"],
        [
            ["KPI Name", "Clear, descriptive name"],
            ["Related Objective", "Which objective does this measure?"],
            ["Formula/Calculation", "How is it calculated?"],
            ["Data Source", "Where does the data come from?"],
            ["Target Value", "What is the goal?"],
            ["Threshold (Acceptable)", "What is the minimum acceptable?"],
            ["Measurement Frequency", "Daily / Weekly / Monthly / Quarterly"],
            ["Owner", "Who is responsible for tracking?"],
            ["Last Updated", "Date of last review"]
        ]
    )

    # ========== ACTION PLAN ==========
    add_content_slide(prs, "Your Action Plan: Next Steps", [
        "Week 1-2: Foundation",
        "☐ Review organizational strategy and departmental objectives",
        "☐ Finalize your function's 3-5 SMART objectives",
        "☐ Get approval from your manager",
        "",
        "Week 3-4: KPIs and Cascading",
        "☐ Define 3-5 KPIs for each function objective",
        "☐ Set up tracking mechanisms and dashboards",
        "☐ Begin cascading: Meet 1-on-1 with each team member",
        "☐ Co-create individual objectives and KPIs",
        "",
        "Ongoing:",
        "☐ Track KPIs at defined frequency",
        "☐ Hold regular review meetings",
        "☐ Adjust as needed based on business changes"
    ])

    # ========== SUCCESS FACTORS ==========
    add_content_slide(prs, "Critical Success Factors", [
        "To succeed with objectives and KPIs:",
        "",
        "1. Leadership Commitment",
        "   • Visible support from senior management",
        "",
        "2. Clear Communication",
        "   • Everyone understands the 'why' and 'how'",
        "",
        "3. Data Availability",
        "   • Can actually measure what you define",
        "",
        "4. Regular Reviews",
        "   • Make it part of your rhythm, not a one-time event",
        "",
        "5. Flexibility",
        "   • Adjust when context changes, but maintain stability",
        "",
        "6. Recognition",
        "   • Celebrate achievements and progress"
    ])

    # ========== WORKSHOP SUMMARY ==========
    add_content_slide(prs, "Workshop Summary: Key Takeaways", [
        "Day 1 - SMART Objectives:",
        "✓ Objectives define what you want to achieve",
        "✓ SMART framework ensures clarity and effectiveness",
        "✓ Good objectives align with organizational strategy",
        "",
        "Day 2 - KPIs & Cascading:",
        "✓ KPIs are quantifiable measures of objective achievement",
        "✓ Different types of KPIs serve different purposes",
        "✓ Cascading creates alignment from company to individual",
        "✓ Regular tracking and review drives continuous improvement",
        "",
        "Remember: This is a journey, not a destination!",
        "Start small, learn, and iterate."
    ])

    # ========== Q&A ==========
    add_title_slide(prs, "Questions & Discussion", "Thank you for your participation!")

    # ========== RESOURCES ==========
    add_content_slide(prs, "Additional Resources", [
        "Recommended Reading:",
        "• 'Measure What Matters' by John Doerr (OKRs)",
        "• 'The Balanced Scorecard' by Kaplan & Norton",
        "• 'High Output Management' by Andy Grove",
        "",
        "Online Resources:",
        "• Google's Guide to OKRs: rework.withgoogle.com/guides/set-goals-with-okrs",
        "• KPI Library: kpilibrary.com",
        "• Atlassian Goal Setting Guide: atlassian.com/team-playbook/plays/goals-signals-measures",
        "",
        "Tools:",
        "• Free templates available on Google Sheets",
        "• Explore OKR software trials (Weekdone, Perdoo)",
        "• Dashboard tools: Grafana (open source)"
    ])

    # ========== CONTACT ==========
    add_title_slide(prs, "Stay Connected", "Keep practicing and don't hesitate to reach out for support!")

    return prs


def main():
    print("Creating comprehensive KPIs and SMART Objectives presentation...")
    prs = create_presentation()

    output_file = "/home/user/Azka-Company-Workshop/How_to_Write_Effective_KPIs_and_SMART_Objectives.pptx"
    prs.save(output_file)

    print(f"✅ Presentation created successfully!")
    print(f"📁 Location: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print("\nThe presentation includes:")
    print("  • Comprehensive 2-day workshop structure")
    print("  • Detailed SMART framework explanation")
    print("  • Multiple tables, examples, and templates")
    print("  • Technology sector-specific examples")
    print("  • Cascading framework and process")
    print("  • Practical exercises for participants")
    print("  • Charts and visual elements")


if __name__ == "__main__":
    main()
